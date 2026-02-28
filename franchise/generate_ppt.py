#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ONE桌遊加盟說明簡報生成器
使用 python-pptx 創建 PowerPoint 檔案
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # 創建簡報
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # 定義顏色
    PRIMARY_COLOR = RGBColor(0, 245, 255)  # 青色
    SECONDARY_COLOR = RGBColor(255, 0, 110)  # 粉紅色
    BG_COLOR = RGBColor(10, 14, 39)  # 深藍色
    
    # ================== Slide 1: 封面 ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(12), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "ONE桌遊"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(96)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    title_p.alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(12), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "加盟說明簡報"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(48)
    subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # 描述
    desc_box = slide.shapes.add_textbox(Inches(2), Inches(6), Inches(12), Inches(0.6))
    desc_frame = desc_box.text_frame
    desc_frame.text = "全台最大自助桌遊連鎖品牌"
    desc_p = desc_frame.paragraphs[0]
    desc_p.font.size = Pt(36)
    desc_p.font.color.rgb = PRIMARY_COLOR
    desc_p.alignment = PP_ALIGN.CENTER
    
    # 公司名
    company_box = slide.shapes.add_textbox(Inches(2), Inches(7.5), Inches(12), Inches(0.4))
    company_frame = company_box.text_frame
    company_frame.text = "遊戲家資訊科技有限公司 出品"
    company_p = company_frame.paragraphs[0]
    company_p.font.size = Pt(18)
    company_p.font.color.rgb = RGBColor(180, 180, 180)
    company_p.alignment = PP_ALIGN.CENTER
    
    # ================== Slide 2: 品牌介紹 ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "品牌介紹"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    title_p.alignment = PP_ALIGN.CENTER
    
    # 內容
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(13), Inches(2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.paragraphs[0]
    p.text = "全台最大自助桌遊連鎖品牌"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    p = content_frame.add_paragraph()
    p.text = "ONE桌遊成立於 2023 年，短短 3 年內快速展店至 105 家，"
    p.text += "運用 AI 智慧運營技術，打造 24 小時無人化自助空間。"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(220, 220, 220)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(20)
    
    # 統計數據
    stats = [
        ("105+", "門市數量"),
        ("10萬", "會員數量"),
        ("24/7", "無人化營運")
    ]
    
    for i, (number, label) in enumerate(stats):
        x_pos = 2 + (i * 4)
        
        # 數字
        num_box = slide.shapes.add_textbox(Inches(x_pos), Inches(5), Inches(3.5), Inches(1))
        num_frame = num_box.text_frame
        num_frame.text = number
        num_p = num_frame.paragraphs[0]
        num_p.font.size = Pt(48)
        num_p.font.bold = True
        num_p.font.color.rgb = PRIMARY_COLOR
        num_p.alignment = PP_ALIGN.CENTER
        
        # 標籤
        label_box = slide.shapes.add_textbox(Inches(x_pos), Inches(6), Inches(3.5), Inches(0.5))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_p = label_frame.paragraphs[0]
        label_p.font.size = Pt(20)
        label_p.font.color.rgb = RGBColor(200, 200, 200)
        label_p.alignment = PP_ALIGN.CENTER
    
    # ================== Slide 3: 市場優勢 ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "市場優勢"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    title_p.alignment = PP_ALIGN.CENTER
    
    # 左側：獨特定位
    left_title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(0.6))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = "🎯 獨特定位"
    left_title_p = left_title_frame.paragraphs[0]
    left_title_p.font.size = Pt(28)
    left_title_p.font.bold = True
    left_title_p.font.color.rgb = PRIMARY_COLOR
    
    left_content_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(7), Inches(4))
    left_content_frame = left_content_box.text_frame
    left_content_frame.word_wrap = True
    
    left_points = [
        "全台首創 24 小時自助麻將館",
        "手機 APP 預約 + AI 智慧門禁",
        "無需櫃檯人員，降低人力成本",
        "重新定義現代桌遊休閒體驗"
    ]
    
    for i, point in enumerate(left_points):
        if i == 0:
            p = left_content_frame.paragraphs[0]
        else:
            p = left_content_frame.add_paragraph()
        p.text = f"✅ {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.space_before = Pt(10)
    
    # 右側：市場需求
    right_title_box = slide.shapes.add_textbox(Inches(8.5), Inches(2), Inches(7), Inches(0.6))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = "📈 市場需求"
    right_title_p = right_title_frame.paragraphs[0]
    right_title_p.font.size = Pt(28)
    right_title_p.font.bold = True
    right_title_p.font.color.rgb = PRIMARY_COLOR
    
    right_content_box = slide.shapes.add_textbox(Inches(8.5), Inches(2.8), Inches(7), Inches(4))
    right_content_frame = right_content_box.text_frame
    right_content_frame.word_wrap = True
    
    right_points = [
        "麻將文化深植台灣社會",
        "年輕族群接受度高",
        "自助服務符合現代消費習慣",
        "疫後聚會需求強勁復甦"
    ]
    
    for i, point in enumerate(right_points):
        if i == 0:
            p = right_content_frame.paragraphs[0]
        else:
            p = right_content_frame.add_paragraph()
        p.text = f"✅ {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.space_before = Pt(10)
    
    # ================== Slide 4: 加盟六大優勢 ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "加盟六大優勢"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    title_p.alignment = PP_ALIGN.CENTER
    
    # 六大優勢
    advantages = [
        ("🎯 AI 智能選址服務", "數據分析最佳地點，降低選址風險"),
        ("🏗️ 裝潢設計規劃", "30-45 天快速開店，標準化流程"),
        ("📱 智能營運管理系統", "手機遠端監控，隨時掌握店況"),
        ("📚 完整教育訓練", "SOP 標準化流程，快速上手"),
        ("🔧 持續技術支援", "24/7 線上客服，即時解決問題"),
        ("📢 品牌行銷資源", "共享 10 萬會員流量，快速導客")
    ]
    
    for i, (title, desc) in enumerate(advantages):
        row = i // 2
        col = i % 2
        x_pos = 1.5 + (col * 7)
        y_pos = 2 + (row * 1.8)
        
        # 標題
        title_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(6.5), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(22)
        title_p.font.bold = True
        title_p.font.color.rgb = PRIMARY_COLOR
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.5), Inches(6.5), Inches(0.6))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_frame.text = desc
        desc_p = desc_frame.paragraphs[0]
        desc_p.font.size = Pt(16)
        desc_p.font.color.rgb = RGBColor(200, 200, 200)
    
    # ================== Slide 5: 投資回報分析 ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "投資回報分析"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    title_p.alignment = PP_ALIGN.CENTER
    
    # 回本期 & 營收
    roi_stats = [
        ("12-18", "個月回本期"),
        ("15-50", "萬/月營收")
    ]
    
    for i, (number, label) in enumerate(roi_stats):
        x_pos = 4 + (i * 4.5)
        
        # 數字
        num_box = slide.shapes.add_textbox(Inches(x_pos), Inches(2), Inches(4), Inches(1))
        num_frame = num_box.text_frame
        num_frame.text = number
        num_p = num_frame.paragraphs[0]
        num_p.font.size = Pt(60)
        num_p.font.bold = True
        num_p.font.color.rgb = PRIMARY_COLOR
        num_p.alignment = PP_ALIGN.CENTER
        
        # 標籤
        label_box = slide.shapes.add_textbox(Inches(x_pos), Inches(3), Inches(4), Inches(0.5))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_p = label_frame.paragraphs[0]
        label_p.font.size = Pt(24)
        label_p.font.color.rgb = RGBColor(200, 200, 200)
        label_p.alignment = PP_ALIGN.CENTER
    
    # 投資明細
    investment_details = [
        ("💵 初期投資", [
            "加盟金：視區域而定",
            "裝潢費用：80-150 萬",
            "設備費用：50-80 萬",
            "押金雜支：20-30 萬"
        ]),
        ("📊 每月支出", [
            "租金：3-8 萬",
            "水電費：1-2 萬",
            "系統費用：8,000-12,000",
            "雜支：5,000-10,000"
        ])
    ]
    
    for i, (category, items) in enumerate(investment_details):
        x_pos = 1.5 + (i * 7)
        
        # 類別標題
        cat_title_box = slide.shapes.add_textbox(Inches(x_pos), Inches(4.5), Inches(6.5), Inches(0.5))
        cat_title_frame = cat_title_box.text_frame
        cat_title_frame.text = category
        cat_title_p = cat_title_frame.paragraphs[0]
        cat_title_p.font.size = Pt(24)
        cat_title_p.font.bold = True
        cat_title_p.font.color.rgb = PRIMARY_COLOR
        
        # 項目列表
        items_box = slide.shapes.add_textbox(Inches(x_pos), Inches(5.2), Inches(6.5), Inches(2.5))
        items_frame = items_box.text_frame
        items_frame.word_wrap = True
        
        for j, item in enumerate(items):
            if j == 0:
                p = items_frame.paragraphs[0]
            else:
                p = items_frame.add_paragraph()
            p.text = f"✅ {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(220, 220, 220)
            p.space_before = Pt(8)
    
    # ================== Slide 6: 成功案例 ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "成功案例"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    title_p.alignment = PP_ALIGN.CENTER
    
    # 主標題
    main_title_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(12), Inches(0.8))
    main_title_frame = main_title_box.text_frame
    main_title_frame.text = "3 年內快速展店至 105 家"
    main_title_p = main_title_frame.paragraphs[0]
    main_title_p.font.size = Pt(36)
    main_title_p.font.bold = True
    main_title_p.font.color.rgb = PRIMARY_COLOR
    main_title_p.alignment = PP_ALIGN.CENTER
    
    # 時間軸
    timeline = [
        ("2023/7", "第一家店開幕"),
        ("2024", "擴展至 50 家店"),
        ("2025", "突破 100 家店")
    ]
    
    for i, (date, milestone) in enumerate(timeline):
        x_pos = 2.5 + (i * 4)
        
        # 日期
        date_box = slide.shapes.add_textbox(Inches(x_pos), Inches(3.5), Inches(3.5), Inches(0.8))
        date_frame = date_box.text_frame
        date_frame.text = date
        date_p = date_frame.paragraphs[0]
        date_p.font.size = Pt(40)
        date_p.font.bold = True
        date_p.font.color.rgb = PRIMARY_COLOR
        date_p.alignment = PP_ALIGN.CENTER
        
        # 里程碑
        milestone_box = slide.shapes.add_textbox(Inches(x_pos), Inches(4.3), Inches(3.5), Inches(0.6))
        milestone_frame = milestone_box.text_frame
        milestone_frame.text = milestone
        milestone_p = milestone_frame.paragraphs[0]
        milestone_p.font.size = Pt(20)
        milestone_p.font.color.rgb = RGBColor(200, 200, 200)
        milestone_p.alignment = PP_ALIGN.CENTER
    
    # AI 技術導入
    ai_box = slide.shapes.add_textbox(Inches(3), Inches(6), Inches(10), Inches(1.5))
    ai_frame = ai_box.text_frame
    
    p = ai_frame.paragraphs[0]
    p.text = "🎉 2025/11 正式導入 AI 技術"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    p = ai_frame.add_paragraph()
    p.text = "全面升級智能化管理，加盟主營運更輕鬆"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(220, 220, 220)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)
    
    # ================== Slide 7: 加盟流程 ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "加盟流程"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    title_p.alignment = PP_ALIGN.CENTER
    
    # 流程步驟
    steps = [
        ("1️⃣ 諮詢評估", "填寫加盟意願書，初步評估"),
        ("2️⃣ 簽約準備", "確認合約內容，繳交訂金"),
        ("3️⃣ 選址裝潢", "AI 選址 + 30-45 天裝潢"),
        ("4️⃣ 教育訓練", "完整 SOP 培訓，系統操作"),
        ("5️⃣ 開幕籌備", "設備安裝、系統測試"),
        ("6️⃣ 正式營運", "開幕活動 + 持續技術支援")
    ]
    
    for i, (step, desc) in enumerate(steps):
        row = i // 2
        col = i % 2
        x_pos = 1.5 + (col * 7)
        y_pos = 2 + (row * 1.8)
        
        # 步驟標題
        step_title_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(6.5), Inches(0.5))
        step_title_frame = step_title_box.text_frame
        step_title_frame.text = step
        step_title_p = step_title_frame.paragraphs[0]
        step_title_p.font.size = Pt(22)
        step_title_p.font.bold = True
        step_title_p.font.color.rgb = PRIMARY_COLOR
        
        # 步驟描述
        step_desc_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.5), Inches(6.5), Inches(0.6))
        step_desc_frame = step_desc_box.text_frame
        step_desc_frame.word_wrap = True
        step_desc_frame.text = desc
        step_desc_p = step_desc_frame.paragraphs[0]
        step_desc_p.font.size = Pt(18)
        step_desc_p.font.color.rgb = RGBColor(200, 200, 200)
    
    # ================== Slide 8: 為什麼選擇 ONE桌遊？ ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "為什麼選擇 ONE桌遊？"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    title_p.alignment = PP_ALIGN.CENTER
    
    # 四大優勢
    reasons = [
        ("✅ 品牌優勢", ["全台最大連鎖品牌", "105 家店成功驗證", "10 萬會員基礎"]),
        ("✅ 技術優勢", ["AI 智慧運營系統", "24/7 無人化管理", "遠端監控降低成本"]),
        ("✅ 獲利優勢", ["12-18 個月回本", "月營收 15-50 萬", "低人力成本"]),
        ("✅ 支援優勢", ["完整教育訓練", "24/7 技術支援", "持續系統升級"])
    ]
    
    for i, (title, points) in enumerate(reasons):
        row = i // 2
        col = i % 2
        x_pos = 1.5 + (col * 7)
        y_pos = 2.2 + (row * 3)
        
        # 標題
        reason_title_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(6.5), Inches(0.5))
        reason_title_frame = reason_title_box.text_frame
        reason_title_frame.text = title
        reason_title_p = reason_title_frame.paragraphs[0]
        reason_title_p.font.size = Pt(24)
        reason_title_p.font.bold = True
        reason_title_p.font.color.rgb = PRIMARY_COLOR
        
        # 項目列表
        points_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.6), Inches(6.5), Inches(1.8))
        points_frame = points_box.text_frame
        points_frame.word_wrap = True
        
        for j, point in enumerate(points):
            if j == 0:
                p = points_frame.paragraphs[0]
            else:
                p = points_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(220, 220, 220)
            p.space_before = Pt(6)
    
    # ================== Slide 9: 聯絡我們 ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    # 主標題
    main_title_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(12), Inches(1))
    main_title_frame = main_title_box.text_frame
    main_title_frame.text = "立即加入 ONE桌遊"
    main_title_p = main_title_frame.paragraphs[0]
    main_title_p.font.size = Pt(60)
    main_title_p.font.bold = True
    main_title_p.font.color.rgb = PRIMARY_COLOR
    main_title_p.alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(12), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    
    p = subtitle_frame.paragraphs[0]
    p.text = "開啟您的創業之路"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    p = subtitle_frame.add_paragraph()
    p.text = "全台最大自助桌遊連鎖品牌 與您一起打造成功事業"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(220, 220, 220)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(15)
    
    # 聯絡資訊
    contact_box = slide.shapes.add_textbox(Inches(4), Inches(5.5), Inches(8), Inches(2))
    contact_frame = contact_box.text_frame
    
    contact_info = [
        "聯絡資訊",
        "",
        "📞 電話：待提供",
        "📧 Email：待提供",
        "🌐 官網：https://ken0900101318-beep.github.io/touchmo-website/"
    ]
    
    for i, line in enumerate(contact_info):
        if i == 0:
            p = contact_frame.paragraphs[0]
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
        else:
            p = contact_frame.add_paragraph()
            if i == 1:
                p.font.size = Pt(10)
            else:
                p.font.size = Pt(22)
                p.font.color.rgb = RGBColor(220, 220, 220)
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        if i > 1:
            p.space_before = Pt(10)
    
    # 儲存簡報
    prs.save('ONE桌遊加盟說明簡報.pptx')
    print("✅ PowerPoint 簡報已成功創建：ONE桌遊加盟說明簡報.pptx")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ 錯誤：缺少 python-pptx 套件")
        print("請執行以下指令安裝：")
        print("pip3 install python-pptx")
    except Exception as e:
        print(f"❌ 創建簡報時發生錯誤：{e}")
