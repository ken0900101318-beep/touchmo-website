#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ONE桌遊加盟說明簡報生成器（明亮版 - 強調高報酬、法規、市場供需）
使用 python-pptx 創建 PowerPoint 檔案
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # 創建簡報
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # 定義顏色（明亮版）
    PRIMARY_COLOR = RGBColor(255, 107, 53)  # 橘色
    SECONDARY_COLOR = RGBColor(247, 147, 30)  # 金橙色
    SUCCESS_COLOR = RGBColor(46, 204, 113)  # 綠色
    BG_COLOR = RGBColor(255, 255, 255)  # 白色
    TEXT_COLOR = RGBColor(51, 51, 51)  # 深灰色
    
    # ================== Slide 1: 封面 ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 250, 245)
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(12), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "ONE桌遊"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(96)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    title_p.alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(12), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "加盟說明簡報"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(48)
    subtitle_p.font.color.rgb = TEXT_COLOR
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # 高報酬標語
    profit_box = slide.shapes.add_textbox(Inches(3), Inches(5.2), Inches(10), Inches(1))
    profit_frame = profit_box.text_frame
    profit_frame.text = "💰 高報酬投資機會"
    profit_p = profit_frame.paragraphs[0]
    profit_p.font.size = Pt(42)
    profit_p.font.bold = True
    profit_p.font.color.rgb = SECONDARY_COLOR
    profit_p.alignment = PP_ALIGN.CENTER
    
    desc_box = slide.shapes.add_textbox(Inches(3), Inches(6), Inches(10), Inches(0.6))
    desc_frame = desc_box.text_frame
    desc_frame.text = "12-18 個月回本 · 月營收 15-50 萬"
    desc_p = desc_frame.paragraphs[0]
    desc_p.font.size = Pt(30)
    desc_p.font.color.rgb = RGBColor(85, 85, 85)
    desc_p.alignment = PP_ALIGN.CENTER
    
    # ================== Slide 2: 市場供需分析 ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "麻將市場供需分析"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    title_p.alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(1.6), Inches(12), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "🀄 台灣麻將文化深厚，市場需求龐大"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(32)
    subtitle_p.font.color.rgb = PRIMARY_COLOR
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # 左側：市場需求
    left_title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(7), Inches(0.6))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = "📈 市場需求"
    left_title_p = left_title_frame.paragraphs[0]
    left_title_p.font.size = Pt(28)
    left_title_p.font.bold = True
    left_title_p.font.color.rgb = PRIMARY_COLOR
    
    left_content_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(7), Inches(4))
    left_content_frame = left_content_box.text_frame
    left_content_frame.word_wrap = True
    
    left_points = [
        "台灣 2,300 萬人口，超過 500 萬人會打麻將",
        "麻將文化深植三代人（50-80 歲主力）",
        "年輕族群接受度逐年提升（16-40 歲）",
        "朋友聚會、家庭聚餐的首選娛樂",
        "疫後聚會需求強勁復甦"
    ]
    
    for i, point in enumerate(left_points):
        if i == 0:
            p = left_content_frame.paragraphs[0]
        else:
            p = left_content_frame.add_paragraph()
        p.text = f"✅ {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
    
    # 右側：市場供給不足
    right_title_box = slide.shapes.add_textbox(Inches(8.5), Inches(2.5), Inches(7), Inches(0.6))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = "📉 市場供給不足"
    right_title_p = right_title_frame.paragraphs[0]
    right_title_p.font.size = Pt(28)
    right_title_p.font.bold = True
    right_title_p.font.color.rgb = PRIMARY_COLOR
    
    right_content_box = slide.shapes.add_textbox(Inches(8.5), Inches(3.2), Inches(7), Inches(4))
    right_content_frame = right_content_box.text_frame
    right_content_frame.word_wrap = True
    
    right_points = [
        "傳統麻將館老舊、環境差",
        "需要 4 人湊局，門檻高",
        "營業時間受限（多數僅白天營業）",
        "年輕人不敢進入傳統麻將館",
        "市場缺口極大，等待創新品牌填補"
    ]
    
    for i, point in enumerate(right_points):
        if i == 0:
            p = right_content_frame.paragraphs[0]
        else:
            p = right_content_frame.add_paragraph()
        p.text = f"✅ {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
    
    # 商機總結
    summary_box = slide.shapes.add_textbox(Inches(2.5), Inches(7.5), Inches(11), Inches(0.8))
    summary_frame = summary_box.text_frame
    summary_frame.text = "💡 商機：龐大需求 × 供給不足 = 高獲利空間"
    summary_p = summary_frame.paragraphs[0]
    summary_p.font.size = Pt(28)
    summary_p.font.bold = True
    summary_p.font.color.rgb = SECONDARY_COLOR
    summary_p.alignment = PP_ALIGN.CENTER
    
    # ================== Slide 3: 法規合法說明 ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "法規合法說明"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    title_p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(1.6), Inches(12), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "✅ 完全合法經營，符合政府法規"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(32)
    subtitle_p.font.color.rgb = SUCCESS_COLOR
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # 四個法規要點
    legal_points = [
        ("📋 營業登記", [
            "營業項目：桌遊館、麻將館",
            "合法商業登記",
            "依法繳稅、開立發票",
            "符合消防、建管法規"
        ]),
        ("🚫 禁止賭博", [
            "僅提供場地租賃服務",
            "不涉及金錢賭博",
            "館內明確張貼「禁止賭博」告示",
            "AI 監控系統預防違法行為"
        ]),
        ("🏢 使用執照", [
            "商業區、住宅區皆可（依縣市而定）",
            "需符合使用分區規定",
            "遠離國中小學 200 公尺",
            "我們協助選址，確保合法"
        ]),
        ("🛡️ 法律顧問", [
            "專業律師團隊協助",
            "協助處理法規問題",
            "定期法規更新通知",
            "加盟主無後顧之憂"
        ])
    ]
    
    for i, (category, items) in enumerate(legal_points):
        row = i // 2
        col = i % 2
        x_pos = 1 + (col * 7.5)
        y_pos = 2.5 + (row * 2.5)
        
        # 類別標題
        cat_title_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(7), Inches(0.5))
        cat_title_frame = cat_title_box.text_frame
        cat_title_frame.text = category
        cat_title_p = cat_title_frame.paragraphs[0]
        cat_title_p.font.size = Pt(22)
        cat_title_p.font.bold = True
        cat_title_p.font.color.rgb = PRIMARY_COLOR
        
        # 項目列表
        items_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.6), Inches(7), Inches(1.8))
        items_frame = items_box.text_frame
        items_frame.word_wrap = True
        
        for j, item in enumerate(items):
            if j == 0:
                p = items_frame.paragraphs[0]
            else:
                p = items_frame.add_paragraph()
            p.text = f"✅ {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(5)
    
    # 總結
    summary_box = slide.shapes.add_textbox(Inches(2), Inches(7.8), Inches(12), Inches(0.6))
    summary_frame = summary_box.text_frame
    summary_frame.text = "✅ 我們的 105 家店全部合法經營，政府輔導、銀行認可"
    summary_p = summary_frame.paragraphs[0]
    summary_p.font.size = Pt(24)
    summary_p.font.bold = True
    summary_p.font.color.rgb = SUCCESS_COLOR
    summary_p.alignment = PP_ALIGN.CENTER
    
    # ================== Slide 4: 高報酬投資回報 ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "高報酬投資回報分析"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    title_p.alignment = PP_ALIGN.CENTER
    
    # 回本期 & 營收
    roi_stats = [
        ("12-18", "個月回本期"),
        ("15-50", "萬/月營收")
    ]
    
    for i, (number, label) in enumerate(roi_stats):
        x_pos = 4 + (i * 4.5)
        
        num_box = slide.shapes.add_textbox(Inches(x_pos), Inches(2), Inches(4), Inches(1))
        num_frame = num_box.text_frame
        num_frame.text = number
        num_p = num_frame.paragraphs[0]
        num_p.font.size = Pt(70)
        num_p.font.bold = True
        num_p.font.color.rgb = SECONDARY_COLOR
        num_p.alignment = PP_ALIGN.CENTER
        
        label_box = slide.shapes.add_textbox(Inches(x_pos), Inches(3), Inches(4), Inches(0.5))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_p = label_frame.paragraphs[0]
        label_p.font.size = Pt(24)
        label_p.font.color.rgb = TEXT_COLOR
        label_p.alignment = PP_ALIGN.CENTER
    
    # 投資明細
    investment_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(6.5), Inches(3))
    investment_frame = investment_box.text_frame
    
    p = investment_frame.paragraphs[0]
    p.text = "💵 初期投資（150-280 萬）"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    items = [
        "加盟金：30-50 萬",
        "裝潢費用：80-150 萬",
        "設備費用：50-80 萬",
        "押金雜支：20-30 萬"
    ]
    
    for item in items:
        p = investment_frame.add_paragraph()
        p.text = f"✅ {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
    
    # 每月支出
    expense_box = slide.shapes.add_textbox(Inches(8.5), Inches(4.2), Inches(6.5), Inches(3))
    expense_frame = expense_box.text_frame
    
    p = expense_frame.paragraphs[0]
    p.text = "📊 每月支出（5-12 萬）"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    items = [
        "租金：3-8 萬",
        "水電費：1-2 萬",
        "系統費用：8,000-12,000",
        "雜支維護：5,000-10,000"
    ]
    
    for item in items:
        p = expense_frame.add_paragraph()
        p.text = f"✅ {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
    
    # 獲利試算
    profit_box = slide.shapes.add_textbox(Inches(2), Inches(7.5), Inches(12), Inches(0.9))
    profit_frame = profit_box.text_frame
    
    p = profit_frame.paragraphs[0]
    p.text = "💰 獲利試算：月營收 25 萬 - 月支出 8 萬 = 月淨利 17 萬"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    p = profit_frame.add_paragraph()
    p.text = "投資 200 萬 ÷ 月淨利 17 萬 = 12 個月回本"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # 繼續創建其他投影片...（簡化版）
    # 儲存簡報
    prs.save('ONE桌遊加盟說明簡報_明亮版.pptx')
    print("✅ PowerPoint 簡報（明亮版）已成功創建：ONE桌遊加盟說明簡報_明亮版.pptx")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ 錯誤：缺少 python-pptx 套件")
        print("請執行以下指令安裝：")
        print("pip3 install python-pptx")
    except Exception as e:
        print(f"❌ 創建簡報時發生錯誤：{e}")
